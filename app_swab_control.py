
import io
import re
import calendar

import numpy as np
import pandas as pd
import plotly.express as px
import streamlit as st

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


st.set_page_config(
    page_title="SWAB Lote X - Servicio y Potencial",
    page_icon="🛢️",
    layout="wide"
)


# ============================================================
# LISTAS FIJAS
# ============================================================

CONVERTIDOS_2024 = [
    "EA11934D", "EA 1868", "EA 8291", "EA 8561", "EA 8547",
    "EA11873D", "EA11054", "EA 8569", "AA11122D", "EA11204D",
    "EA 1625", "EA 8532", "EA11703D", "EA11591D", "AA 6316",
    "EA 8593", "EA 8506", "AA11976D", "AA 5804", "EA 8209",
    "EA 8773", "AA11958D", "AA11214D", "EA11968D", "EA11221D",
    "EA 5997", "EA 8523", "EA 8783", "EA 5718", "EA 8317",
    "AA 8309D", "AA 6089", "EA11669D", "EA 1763", "EA 8961D",
    "EA11708D", "AA11408D", "EA 8901", "EA 7991", "EA11033",
    "EA11931D", "EA 8716", "EA11748D", "EA 9472", "EA11194",
    "EA11656D", "AA11342D", "EA 8687", "EA11431D"
]

CONVERTIDOS_2025 = [
    "AA11062D", "EA11049D", "EA 8204", "EA 8043", "AA11118D",
    "EA11918D", "AA11702D", "AA11396D", "EA11929D", "EA 8544",
    "EA 2324", "EA 9241", "EA11224D", "EA11884D", "AA 8149D",
    "AA11363D", "AA 6274", "EA 5998", "EA 8088", "EA 8902",
    "EA 8724", "EA11813D", "EA 8759", "EA 7516", "AA 5931",
    "AA11801D"
]

CONVERTIDOS_2026 = [
    "EA 9836", "EA 8417", "EA 9538", "EA 5811", "EA 8672",
    "EA11444", "AA11172D", "EA 8586", "AA   47", "AA11762",
    "EA11609D", "AA11144D", "EA 9041", "EA11388D", "AA11481D",
    "AA11327D", "AA  106", "EA8819", "EA8816D", "EA8662D"
]

CANDIDATOS_ATA = [
    "AA37", "AA54", "AA76", "AA112", "AA1577", "AA1598", "AA1599", "AA1633",
    "AA1661", "AA1847", "AA1930", "AA5631", "AA5707", "AA5861", "AA5926",
    "AA5971", "AA6192", "AA6338", "AA6342", "AA6372", "AA6423", "AA6454",
    "AA6517", "AA6646", "AA6762", "AA7201", "AA9154", "AA9329", "AA9364",
    "AA10013", "EA216", "EA264", "EA364", "EA440", "EA741", "EA771", "EA876",
    "EA888", "EA987", "EA1054", "EA1081", "EA1161", "EA1167", "EA1233",
    "EA1302", "EA1506", "EA1511", "EA1513", "EA1581", "EA1630", "EA1885",
    "EA2067", "EA2249", "EA2254", "EA2256", "EA2304", "EA2372", "EA2389",
    "EA2403", "EA5682D", "EA5694", "EA5739", "EA5766", "EA5868", "EA5874",
    "EA5914", "EA5921", "EA5957", "EA6130", "EA6237", "EA6918", "EA7027",
    "EA7158", "EA8574", "EA9242", "EA9251", "EA9287", "EA9409", "EA9417",
    "EA9491", "EA9668", "EA9752", "EA9779", "EA11128", "PB47", "PB232",
    "PE171", "PT4-3"
]

MESES = {
    1: "Enero", 2: "Febrero", 3: "Marzo", 4: "Abril",
    5: "Mayo", 6: "Junio", 7: "Julio", 8: "Agosto",
    9: "Septiembre", 10: "Octubre", 11: "Noviembre", 12: "Diciembre"
}

LABELS = {
    "POZO": "Pozo",
    "BATERIA": "Batería",
    "CLASIFICACION": "Clasificación",
    "ANIO_CONVERSION": "Año conversión",
    "ESTADO_DESPLAZAMIENTO": "Estado",
    "PRIORIDAD_REVISION": "Prioridad",
    "INTERV_BASE": "Intervenciones base promedio",
    "INTERV_ACTUAL": "Intervenciones mes objetivo",
    "INTERV_BASE_NO_REALIZADAS": "Intervenciones base no realizadas",
    "PRCR_BASE": "Producción de petróleo base",
    "PRCR_ACTUAL": "Producción de petróleo actual",
    "PRCR_BASE_NO_REALIZADO": "Petróleo base no realizado",
    "PRAG_BASE": "Producción de agua base",
    "PRAG_ACTUAL": "Producción de agua actual",
    "PRAG_BASE_NO_REALIZADO": "Agua base no realizada",
    "PRCR_MES_ANTERIOR_OBJETIVO": "Petróleo mes anterior al objetivo",
    "PRODUCCION_PETROLEO_DEJADA_MES_ANTERIOR": "Petróleo dejado vs mes anterior",
    "VAR_PETROLEO_VS_MES_ANTERIOR": "Variación petróleo vs mes anterior",

    "ULTIMO_MES_INTERVENIDO": "Último mes intervenido",
    "ULTIMA_FECHA_INTERVENCION": "Última fecha de intervención",
    "DIAS_ULTIMO_MES_INTERVENIDO": "Días último mes intervenido",
    "PRCR_ULTIMO_MES_INTERVENIDO": "Petróleo último mes intervenido",
    "PRAG_ULTIMO_MES_INTERVENIDO": "Agua último mes intervenido",
    "INTERV_ULTIMO_MES_INTERVENIDO": "Intervenciones último mes intervenido",
    "POTENCIAL_ULTIMO_MES_INTERVENIDO_BOPD": "Potencial último mes intervenido bopd",

    "MES_1_MES_ANTES_ULTIMA_INTERVENCION": "1 mes antes de última intervención",
    "DIAS_1_MES_ANTES_ULTIMA_INTERVENCION": "Días 1 mes antes",
    "PRCR_1_MES_ANTES_ULTIMA_INTERVENCION": "Petróleo 1 mes antes",
    "PRAG_1_MES_ANTES_ULTIMA_INTERVENCION": "Agua 1 mes antes",
    "INTERV_1_MES_ANTES_ULTIMA_INTERVENCION": "Intervenciones 1 mes antes",
    "POTENCIAL_1_MES_ANTES_ULTIMA_INTERVENCION_BOPD": "Potencial 1 mes antes bopd",

    "MES_2_MES_ANTES_ULTIMA_INTERVENCION": "2 meses antes de última intervención",
    "DIAS_2_MES_ANTES_ULTIMA_INTERVENCION": "Días 2 meses antes",
    "PRCR_2_MES_ANTES_ULTIMA_INTERVENCION": "Petróleo 2 meses antes",
    "PRAG_2_MES_ANTES_ULTIMA_INTERVENCION": "Agua 2 meses antes",
    "INTERV_2_MES_ANTES_ULTIMA_INTERVENCION": "Intervenciones 2 meses antes",
    "POTENCIAL_2_MES_ANTES_ULTIMA_INTERVENCION_BOPD": "Potencial 2 meses antes bopd",

    "POTENCIAL_PROMEDIO_3_MESES_INTERVENCION_BOPD": "Potencial promedio 3 meses bopd",
    "ULTIMA_FECHA_CON_PRODUCCION_PETROLEO": "Última fecha con producción de petróleo",
    "PRODUCCION_PETROLEO_ULTIMA_FECHA": "Petróleo en última fecha con producción",
    "PRODUCCION_AGUA_ULTIMA_FECHA": "Agua en última fecha con producción",
    "FECHA_PRODUCCION_PETROLEO_ANTERIOR": "Fecha de producción de petróleo anterior",
    "PRODUCCION_PETROLEO_FECHA_ANTERIOR": "Petróleo en fecha anterior",
    "PRODUCCION_AGUA_FECHA_ANTERIOR": "Agua en fecha anterior",
    "ULTIMO_MES_CON_PRODUCCION_PETROLEO": "Último mes con producción de petróleo",
    "DIAS_ULTIMO_MES_CON_PRODUCCION": "Días último mes con producción",
    "PRCR_ULTIMO_MES_CON_PRODUCCION": "Petróleo último mes con producción",
    "PRAG_ULTIMO_MES_CON_PRODUCCION": "Agua último mes con producción",
    "INTERV_ULTIMO_MES_CON_PRODUCCION": "Intervenciones último mes con producción",
    "POTENCIAL_ULTIMO_MES_CON_PRODUCCION_BOPD": "Potencial último mes con producción bopd",

    "INTERVENCIONES": "Intervenciones",
    "PRCR": "Producción de petróleo",
    "PRAG": "Producción de agua",
    "TIPO_SWAB": "Tipo de swab",
    "OIL_POR_INTERV": "Petróleo por intervención",
    "AGUA_POR_INTERV": "Agua por intervención",
    "ESTADO": "Estado",
    "ULTIMA_FECHA_HISTORICA": "Última fecha histórica"
}


# ============================================================
# FUNCIONES BÁSICAS
# ============================================================

def limpiar_pozo(valor):
    if pd.isna(valor):
        return ""
    return re.sub(r"\s+", "", str(valor).strip().upper())


def mostrar_pozo(valor):
    if pd.isna(valor):
        return ""
    return re.sub(r"\s+", " ", str(valor).strip().upper())


def limpiar_texto(valor):
    if pd.isna(valor):
        return ""
    return re.sub(r"\s+", " ", str(valor).strip().upper())


def normalizar_columna(col):
    texto = str(col).strip().upper()
    for a, b in {"Á": "A", "É": "E", "Í": "I", "Ó": "O", "Ú": "U"}.items():
        texto = texto.replace(a, b)
    return re.sub(r"\s+", "_", texto)


def normalizar_tipo_swab(valor):
    texto = limpiar_texto(valor)
    if texto in ["TBG", "TB", "TS", "TUBING", "TUBING SWAB"] or "TBG" in texto or "TUB" in texto:
        return "TS"
    if texto in ["CS", "CSG", "CASING", "CASING SWAB"] or "CS" in texto or "CAS" in texto:
        return "CS"
    if texto == "":
        return "SIN TIPO"
    return texto


CONV_MAP = {}
for p in CONVERTIDOS_2024:
    CONV_MAP[limpiar_pozo(p)] = ("Convertido 2024", 2024)
for p in CONVERTIDOS_2025:
    CONV_MAP[limpiar_pozo(p)] = ("Convertido 2025", 2025)
for p in CONVERTIDOS_2026:
    CONV_MAP[limpiar_pozo(p)] = ("Convertido 2026", 2026)

ATA_SET = {limpiar_pozo(p) for p in CANDIDATOS_ATA}


def clasificar_pozo(pozo_key):
    return CONV_MAP.get(pozo_key, ("Básica", 0))[0]


def anio_conversion(pozo_key):
    return CONV_MAP.get(pozo_key, ("Básica", 0))[1]


def primer_dia_mes(anio, mes):
    return pd.Timestamp(int(anio), int(mes), 1)


def ultimo_dia_mes(anio, mes):
    return primer_dia_mes(anio, mes) + pd.offsets.MonthEnd(0)


def mes_anterior(anio, mes, n=1):
    fecha = pd.Timestamp(int(anio), int(mes), 1) - pd.DateOffset(months=int(n))
    return int(fecha.year), int(fecha.month)


def lista_meses_previos(anio, mes, n):
    return [mes_anterior(anio, mes, i) for i in range(int(n), 0, -1)]


def dias_mes(anio, mes):
    return calendar.monthrange(int(anio), int(mes))[1]


def periodo_texto(anio, mes):
    return f"{MESES[int(mes)]} {int(anio)}"


def ultimo_mes_completo(df):
    fecha_max = df["FECHA"].max()
    fin_mes = fecha_max + pd.offsets.MonthEnd(0)
    if fecha_max.date() < fin_mes.date():
        ref = fecha_max - pd.DateOffset(months=1)
    else:
        ref = fecha_max
    return int(ref.year), int(ref.month)


def vista_tabla(df):
    salida = df.copy()
    salida = salida.rename(columns={c: LABELS.get(c, c) for c in salida.columns})
    salida.columns = [str(c) for c in salida.columns]

    for col in salida.columns:
        if pd.api.types.is_datetime64_any_dtype(salida[col]):
            salida[col] = salida[col].dt.strftime("%Y-%m-%d")
        elif pd.api.types.is_period_dtype(salida[col]):
            salida[col] = salida[col].astype(str)
        elif pd.api.types.is_numeric_dtype(salida[col]):
            salida[col] = salida[col].fillna(0).round(2)
        else:
            salida[col] = salida[col].fillna("").astype(str)

    return salida


def excel_descarga(tablas):
    buffer = io.BytesIO()
    with pd.ExcelWriter(buffer, engine="xlsxwriter") as writer:
        for nombre, tabla in tablas.items():
            if tabla is None:
                continue
            out = vista_tabla(tabla)
            out.to_excel(writer, sheet_name=str(nombre)[:31], index=False)
            ws = writer.sheets[str(nombre)[:31]]
            wb = writer.book
            header = wb.add_format({"bold": True, "bg_color": "#D9EAF7", "border": 1})
            for j, col in enumerate(out.columns):
                ws.write(0, j, col, header)
                ws.set_column(j, j, 20)

    buffer.seek(0)
    return buffer.getvalue()



# ============================================================
# PPT EDITABLE
# ============================================================

PPT_AZUL = RGBColor(31, 78, 121)
PPT_GRIS = RGBColor(89, 89, 89)
PPT_ROJO = RGBColor(192, 0, 0)
PPT_VERDE = RGBColor(112, 173, 71)


def _ppt_text(slide, text, x, y, w, h, size=16, bold=False, color=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(size)
    run.font.bold = bool(bold)
    if color is not None:
        run.font.color.rgb = color
    return box


def _ppt_title(slide, title, subtitle=""):
    _ppt_text(slide, title, 0.45, 0.25, 12.4, 0.45, size=24, bold=True, color=PPT_AZUL)
    if subtitle:
        _ppt_text(slide, subtitle, 0.48, 0.78, 12.0, 0.3, size=10, color=PPT_GRIS)


def _ppt_kpi(slide, label, value, x, y, w=2.0):
    _ppt_text(slide, label, x, y, w, 0.25, size=8, color=PPT_GRIS)
    _ppt_text(slide, value, x, y + 0.24, w, 0.42, size=18, bold=True, color=PPT_AZUL)


def _ppt_table(slide, df, x, y, w, h, max_rows=12, font_size=6):
    if df is None or df.empty:
        _ppt_text(slide, "Sin datos", x, y, w, h, size=12, color=PPT_GRIS)
        return

    data = vista_tabla(df).head(max_rows).copy()
    rows = len(data) + 1
    cols = len(data.columns)
    shape = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h))
    table = shape.table

    for j, col in enumerate(data.columns):
        cell = table.cell(0, j)
        cell.text = str(col)
        cell.text_frame.paragraphs[0].font.size = Pt(font_size)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.color.rgb = PPT_AZUL

    for i, (_, row) in enumerate(data.iterrows(), start=1):
        for j, col in enumerate(data.columns):
            cell = table.cell(i, j)
            val = row[col]
            if isinstance(val, float):
                txt = f"{val:,.2f}"
            else:
                txt = "" if pd.isna(val) else str(val)
            cell.text = txt
            cell.text_frame.paragraphs[0].font.size = Pt(font_size)


def _ppt_bar(slide, title, df, cat_col, val_col, x, y, w, h, series_name="Valor", max_rows=12):
    if df is None or df.empty or cat_col not in df.columns or val_col not in df.columns:
        _ppt_text(slide, "Sin datos para gráfico", x, y, w, h, size=12, color=PPT_GRIS)
        return

    data = df[[cat_col, val_col]].dropna().copy()
    data[val_col] = pd.to_numeric(data[val_col], errors="coerce").fillna(0)
    data = data.sort_values(val_col, ascending=False).head(max_rows)
    if data.empty:
        return

    chart_data = CategoryChartData()
    chart_data.categories = [str(x) for x in data[cat_col].tolist()]
    chart_data.add_series(series_name, [float(x) for x in data[val_col].tolist()])
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(x), Inches(y), Inches(w), Inches(h),
        chart_data
    ).chart
    chart.chart_title.has_text_frame = True
    chart.chart_title.text_frame.text = title
    chart.has_legend = False
    try:
        chart.value_axis.tick_labels.number_format = '#,##0.00'
    except Exception:
        pass


def _ppt_line(slide, title, df, cat_col, val_cols, x, y, w, h):
    if df is None or df.empty or cat_col not in df.columns:
        _ppt_text(slide, "Sin datos para tendencia", x, y, w, h, size=12, color=PPT_GRIS)
        return

    data = df.copy()
    chart_data = CategoryChartData()
    chart_data.categories = [str(x) for x in data[cat_col].tolist()]
    for col in val_cols:
        if col in data.columns:
            chart_data.add_series(LABELS.get(col, col), [float(x) for x in pd.to_numeric(data[col], errors='coerce').fillna(0).tolist()])

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS,
        Inches(x), Inches(y), Inches(w), Inches(h),
        chart_data
    ).chart
    chart.chart_title.has_text_frame = True
    chart.chart_title.text_frame.text = title
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    try:
        chart.value_axis.tick_labels.number_format = '#,##0.00'
    except Exception:
        pass


def _conv2026_detalle(list_conv, actual, potencial):
    conv = list_conv[list_conv["CLASIFICACION"] == "Convertido 2026"].copy()
    pot_cols = [
        "POZO_KEY",
        "ULTIMO_MES_INTERVENIDO",
        "ULTIMA_FECHA_INTERVENCION",
        "POTENCIAL_ULTIMO_MES_INTERVENIDO_BOPD",
        "POTENCIAL_PROMEDIO_3_MESES_INTERVENCION_BOPD",
        "ULTIMA_FECHA_CON_PRODUCCION_PETROLEO",
        "PRODUCCION_PETROLEO_ULTIMA_FECHA",
        "FECHA_PRODUCCION_PETROLEO_ANTERIOR",
        "PRODUCCION_PETROLEO_FECHA_ANTERIOR",
        "ULTIMO_MES_CON_PRODUCCION_PETROLEO",
        "POTENCIAL_ULTIMO_MES_CON_PRODUCCION_BOPD"
    ]
    pot_cols = [c for c in pot_cols if c in potencial.columns]
    conv = conv.merge(potencial[pot_cols], on="POZO_KEY", how="left")
    conv["ESTADO_CONVERTIDO_2026"] = np.where(conv["INTERVENCIONES"].fillna(0) > 0, "Intervenido", "No intervenido")
    conv["PETROLEO_POR_INTERVENCION"] = np.where(
        conv["INTERVENCIONES"].fillna(0) > 0,
        conv["PRCR"].fillna(0) / conv["INTERVENCIONES"].fillna(0),
        0
    )
    return conv.sort_values(["ESTADO_CONVERTIDO_2026", "PRCR"], ascending=[True, False])


def crear_ppt_dashboard(periodo, kpis, res_clase, res_bat, res_tipo, res_tend, desp, pot, list_conv, actual, top_n=20):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    afectados = desp[desp["ESTADO_DESPLAZAMIENTO"].isin(["DEJADO DE HACER", "REDUCIDO"])].copy()
    top_dej = afectados.sort_values("PRCR_BASE_NO_REALIZADO", ascending=False).head(top_n)
    top_pot = pot.sort_values("POTENCIAL_PROMEDIO_3_MESES_INTERVENCION_BOPD", ascending=False).head(top_n)
    conv_det = _conv2026_detalle(list_conv, actual, pot)

    # Slide 1
    slide = prs.slides.add_slide(blank)
    _ppt_title(slide, "SWAB Lote X: resumen ejecutivo", f"Periodo objetivo: {periodo}")
    _ppt_kpi(slide, "Pozos universo", f"{kpis['pozos_universo']:,}", 0.55, 1.25)
    _ppt_kpi(slide, "Intervenidos", f"{kpis['pozos_interv']:,}", 2.55, 1.25)
    _ppt_kpi(slide, "No intervenidos", f"{kpis['pozos_no']:,}", 4.55, 1.25)
    _ppt_kpi(slide, "Intervenciones", f"{kpis['interv']:,}", 6.55, 1.25)
    _ppt_kpi(slide, "Petróleo", f"{kpis['petroleo']:,.2f}", 8.55, 1.25)
    _ppt_kpi(slide, "Agua", f"{kpis['agua']:,.2f}", 10.55, 1.25)
    _ppt_bar(slide, "Producción de petróleo por clasificación", res_clase, "CLASIFICACION", "PRCR", 0.7, 2.25, 5.7, 4.5, "Petróleo")
    _ppt_bar(slide, "Top baterías por petróleo", res_bat, "BATERIA", "PRCR", 6.85, 2.25, 5.7, 4.5, "Petróleo")

    # Slide 2
    slide = prs.slides.add_slide(blank)
    _ppt_title(slide, "Impacto de convertidos 2026", "Detalle del grupo priorizado de 20 pozos")
    conv_interv = int(conv_det["INTERVENCIONES"].fillna(0).gt(0).sum()) if not conv_det.empty else 0
    _ppt_kpi(slide, "Convertidos intervenidos", f"{conv_interv} de 20", 0.55, 1.1)
    _ppt_kpi(slide, "Intervenciones", f"{conv_det['INTERVENCIONES'].sum():,.0f}" if not conv_det.empty else "0", 2.8, 1.1)
    _ppt_kpi(slide, "Petróleo", f"{conv_det['PRCR'].sum():,.2f}" if not conv_det.empty else "0.00", 5.05, 1.1)
    _ppt_kpi(slide, "Agua", f"{conv_det['PRAG'].sum():,.2f}" if not conv_det.empty else "0.00", 7.3, 1.1)
    _ppt_kpi(slide, "Petróleo/interv.", f"{(conv_det['PRCR'].sum()/conv_det['INTERVENCIONES'].sum()):,.2f}" if not conv_det.empty and conv_det['INTERVENCIONES'].sum() > 0 else "0.00", 9.55, 1.1)
    _ppt_bar(slide, "Top convertidos 2026 por petróleo", conv_det, "POZO", "PRCR", 0.6, 2.15, 6.1, 4.7, "Petróleo")
    cols = ["POZO", "BATERIA", "ESTADO_CONVERTIDO_2026", "INTERVENCIONES", "PRCR", "PRAG", "POTENCIAL_PROMEDIO_3_MESES_INTERVENCION_BOPD"]
    cols = [c for c in cols if c in conv_det.columns]
    _ppt_table(slide, conv_det[cols], 6.95, 2.15, 5.8, 4.7, max_rows=12, font_size=6)

    # Slide 3
    slide = prs.slides.add_slide(blank)
    _ppt_title(slide, "Pozos dejados o reducidos", "Candidatos desplazados por priorización operativa")
    _ppt_bar(slide, "Top pozos por petróleo base no realizado", top_dej, "POZO", "PRCR_BASE_NO_REALIZADO", 0.6, 1.25, 6.2, 5.6, "Petróleo base no realizado")
    cols = ["ESTADO_DESPLAZAMIENTO", "POZO", "BATERIA", "CLASIFICACION", "INTERV_BASE", "INTERV_ACTUAL", "PRCR_BASE_NO_REALIZADO", "POTENCIAL_PROMEDIO_3_MESES_INTERVENCION_BOPD"]
    cols = [c for c in cols if c in top_dej.columns]
    _ppt_table(slide, top_dej[cols], 7.0, 1.25, 5.8, 5.6, max_rows=13, font_size=6)

    # Slide 4
    slide = prs.slides.add_slide(blank)
    _ppt_title(slide, "Potencial y trazabilidad de producción", "Potencial calculado con último mes intervenido y meses previos")
    _ppt_bar(slide, "Top potencial promedio 3 meses", top_pot, "POZO", "POTENCIAL_PROMEDIO_3_MESES_INTERVENCION_BOPD", 0.6, 1.25, 6.2, 5.6, "BOPD")
    cols = ["POZO", "BATERIA", "ULTIMO_MES_INTERVENIDO", "PRCR_ULTIMO_MES_INTERVENIDO", "POTENCIAL_ULTIMO_MES_INTERVENIDO_BOPD", "ULTIMA_FECHA_CON_PRODUCCION_PETROLEO", "PRODUCCION_PETROLEO_ULTIMA_FECHA"]
    cols = [c for c in cols if c in top_pot.columns]
    _ppt_table(slide, top_pot[cols], 7.0, 1.25, 5.8, 5.6, max_rows=13, font_size=6)

    # Slide 5
    slide = prs.slides.add_slide(blank)
    _ppt_title(slide, "Tendencia mensual y tipo de swab", "Producción e intervenciones por periodo")
    _ppt_line(slide, "Tendencia mensual petróleo y agua", res_tend, "MES_NOMBRE", ["PRCR", "PRAG"], 0.65, 1.25, 6.0, 5.4)
    _ppt_bar(slide, "Intervenciones por tipo de swab", res_tipo, "TIPO_SWAB", "INTERVENCIONES", 7.0, 1.25, 5.5, 5.4, "Intervenciones")

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()

def aplicar_layout(fig, titulo, altura=520):
    fig.update_layout(
        template="plotly_white",
        title=titulo,
        height=altura,
        margin=dict(l=20, r=20, t=70, b=50),
        legend_title_text=""
    )
    fig.update_yaxes(tickformat=",.2f")

    for tr in fig.data:
        if getattr(tr, "type", "") == "bar":
            if getattr(tr, "text", None) is not None:
                tr.texttemplate = "%{text:,.2f}"
                tr.textposition = "outside"
                tr.cliponaxis = False

    return fig


# ============================================================
# CARGA RÁPIDA
# ============================================================

@st.cache_data(show_spinner=False)
def cargar_datos(bytes_excel):
    xls = pd.ExcelFile(io.BytesIO(bytes_excel))
    hoja = "Datos de Swab" if "Datos de Swab" in xls.sheet_names else xls.sheet_names[0]

    header = pd.read_excel(io.BytesIO(bytes_excel), sheet_name=hoja, nrows=0)
    mapa_cols = {normalizar_columna(c): c for c in header.columns}

    requeridas_norm = ["FECHA", "COD_POZ", "COD_BAT", "UNIDAD", "TSER", "PRCR", "PRAG"]
    faltantes = [c for c in requeridas_norm if c not in mapa_cols]

    if faltantes:
        raise ValueError("Faltan columnas obligatorias: " + ", ".join(faltantes))

    usecols = [mapa_cols[c] for c in requeridas_norm]
    df = pd.read_excel(io.BytesIO(bytes_excel), sheet_name=hoja, usecols=usecols)
    df.columns = [normalizar_columna(c) for c in df.columns]

    df["FECHA"] = pd.to_datetime(df["FECHA"], errors="coerce")
    df = df[df["FECHA"].notna()].copy()

    df["POZO"] = df["COD_POZ"].apply(mostrar_pozo)
    df["POZO_KEY"] = df["COD_POZ"].apply(limpiar_pozo)
    df["BATERIA"] = df["COD_BAT"].apply(limpiar_texto)
    df["UNIDAD"] = df["UNIDAD"].apply(limpiar_texto)
    df["TIPO_SWAB"] = df["TSER"].apply(normalizar_tipo_swab)
    df["PRCR"] = pd.to_numeric(df["PRCR"], errors="coerce").fillna(0)
    df["PRAG"] = pd.to_numeric(df["PRAG"], errors="coerce").fillna(0)

    df = df[df["POZO_KEY"] != ""].copy()

    df["ANIO"] = df["FECHA"].dt.year
    df["MES"] = df["FECHA"].dt.month
    df["MES_NOMBRE"] = df["MES"].map(MESES)
    df["PERIODO_MES"] = df["FECHA"].dt.to_period("M")
    df["CLASIFICACION"] = df["POZO_KEY"].apply(clasificar_pozo)
    df["ANIO_CONVERSION"] = df["POZO_KEY"].apply(anio_conversion)

    df_ata = df[df["POZO_KEY"].isin(ATA_SET)].copy()
    df_swab = df[~df["POZO_KEY"].isin(ATA_SET)].copy()

    return df_swab, df_ata, hoja


def construir_universo(df):
    hist = (
        df.sort_values("FECHA")
        .groupby("POZO_KEY", as_index=False)
        .agg(
            POZO=("POZO", "last"),
            BATERIA=("BATERIA", "last"),
            ULTIMA_FECHA_HISTORICA=("FECHA", "max")
        )
    )

    conv = []
    for key, (clase, anio) in CONV_MAP.items():
        if key not in ATA_SET:
            conv.append({"POZO_KEY": key, "POZO_CONV": key, "CLASIFICACION_CONV": clase, "ANIO_CONVERSION_CONV": anio})
    conv = pd.DataFrame(conv)

    universo = hist.merge(conv, on="POZO_KEY", how="outer")
    universo["POZO"] = universo["POZO"].fillna(universo["POZO_CONV"])
    universo["BATERIA"] = universo["BATERIA"].fillna("SIN BATERIA")
    universo["CLASIFICACION"] = universo["CLASIFICACION_CONV"].fillna("Básica")
    universo["ANIO_CONVERSION"] = universo["ANIO_CONVERSION_CONV"].fillna(0).astype(int)

    return universo[["POZO_KEY", "POZO", "BATERIA", "CLASIFICACION", "ANIO_CONVERSION", "ULTIMA_FECHA_HISTORICA"]]


# ============================================================
# CÁLCULOS
# ============================================================

def resumen_pozo_periodo(data, sufijo, divisor=1):
    if data.empty:
        return pd.DataFrame(columns=["POZO_KEY", f"INTERV_{sufijo}", f"PRCR_{sufijo}", f"PRAG_{sufijo}", f"OIL_INTERV_{sufijo}"])

    res = (
        data.groupby("POZO_KEY", as_index=False)
        .agg(
            INTERV=("FECHA", "count"),
            PRCR=("PRCR", "sum"),
            PRAG=("PRAG", "sum")
        )
    )

    divisor = max(float(divisor), 1.0)
    res[f"INTERV_{sufijo}"] = res["INTERV"] / divisor
    res[f"PRCR_{sufijo}"] = res["PRCR"] / divisor
    res[f"PRAG_{sufijo}"] = res["PRAG"] / divisor
    res[f"OIL_INTERV_{sufijo}"] = np.where(res[f"INTERV_{sufijo}"] > 0, res[f"PRCR_{sufijo}"] / res[f"INTERV_{sufijo}"], 0)

    return res[["POZO_KEY", f"INTERV_{sufijo}", f"PRCR_{sufijo}", f"PRAG_{sufijo}", f"OIL_INTERV_{sufijo}"]]


def calcular_potencial(df, universo, baterias=None, tipos=None, clases=None):
    """
    Calcula potencial usando como referencia la última intervención del pozo.

    Además agrega trazabilidad de producción real de petróleo:
    1. Última fecha con producción de petróleo mayor a cero
    2. Fecha de producción de petróleo anterior
    3. Último mes con producción de petróleo y su potencial BOPD

    Fórmula de potencial:
    Potencial BOPD = producción de petróleo del mes / días calendario del mes
    """
    data = df.copy()

    if baterias:
        data = data[data["BATERIA"].isin(baterias)]
    if tipos:
        data = data[data["TIPO_SWAB"].isin(tipos)]
    if clases:
        data = data[data["CLASIFICACION"].isin(clases)]

    columnas_base_sin_data = {
        "ULTIMO_MES_INTERVENIDO": "Sin intervención",
        "ULTIMA_FECHA_INTERVENCION": pd.NaT,
        "POTENCIAL_ULTIMO_MES_INTERVENIDO_BOPD": 0.0,
        "POTENCIAL_1_MES_ANTES_ULTIMA_INTERVENCION_BOPD": 0.0,
        "POTENCIAL_2_MES_ANTES_ULTIMA_INTERVENCION_BOPD": 0.0,
        "POTENCIAL_PROMEDIO_3_MESES_INTERVENCION_BOPD": 0.0,
        "ULTIMA_FECHA_CON_PRODUCCION_PETROLEO": pd.NaT,
        "PRODUCCION_PETROLEO_ULTIMA_FECHA": 0.0,
        "PRODUCCION_AGUA_ULTIMA_FECHA": 0.0,
        "FECHA_PRODUCCION_PETROLEO_ANTERIOR": pd.NaT,
        "PRODUCCION_PETROLEO_FECHA_ANTERIOR": 0.0,
        "PRODUCCION_AGUA_FECHA_ANTERIOR": 0.0,
        "ULTIMO_MES_CON_PRODUCCION_PETROLEO": "Sin producción",
        "DIAS_ULTIMO_MES_CON_PRODUCCION": 0,
        "PRCR_ULTIMO_MES_CON_PRODUCCION": 0.0,
        "PRAG_ULTIMO_MES_CON_PRODUCCION": 0.0,
        "INTERV_ULTIMO_MES_CON_PRODUCCION": 0,
        "POTENCIAL_ULTIMO_MES_CON_PRODUCCION_BOPD": 0.0
    }

    if data.empty:
        out = universo.copy()
        for col, valor in columnas_base_sin_data.items():
            out[col] = valor
        return out

    data["ANIO_MES"] = data["FECHA"].dt.to_period("M")

    mensual = (
        data.groupby(["POZO_KEY", "ANIO_MES"], as_index=False)
        .agg(
            PRCR_MES=("PRCR", "sum"),
            PRAG_MES=("PRAG", "sum"),
            INTERV_MES=("FECHA", "count"),
            ULTIMA_FECHA_MES=("FECHA", "max")
        )
    )

    idx = mensual.groupby("POZO_KEY")["ANIO_MES"].idxmax()
    ultimo = mensual.loc[idx, ["POZO_KEY", "ANIO_MES", "ULTIMA_FECHA_MES"]].copy()
    ultimo = ultimo.rename(columns={
        "ANIO_MES": "PERIODO_ULTIMA_INTERVENCION",
        "ULTIMA_FECHA_MES": "ULTIMA_FECHA_INTERVENCION"
    })

    registros = []

    for _, row in ultimo.iterrows():
        pozo_key = row["POZO_KEY"]
        periodo_ult = row["PERIODO_ULTIMA_INTERVENCION"]
        anio_ult = int(periodo_ult.year)
        mes_ult = int(periodo_ult.month)

        registro = {
            "POZO_KEY": pozo_key,
            "ULTIMO_MES_INTERVENIDO": periodo_texto(anio_ult, mes_ult),
            "ULTIMA_FECHA_INTERVENCION": row["ULTIMA_FECHA_INTERVENCION"],
        }

        for offset, etiqueta in [
            (0, "ULTIMO_MES_INTERVENIDO"),
            (1, "1_MES_ANTES_ULTIMA_INTERVENCION"),
            (2, "2_MES_ANTES_ULTIMA_INTERVENCION"),
        ]:
            if offset == 0:
                a_ref, m_ref = anio_ult, mes_ult
            else:
                a_ref, m_ref = mes_anterior(anio_ult, mes_ult, offset)

            periodo_ref = pd.Period(f"{a_ref}-{m_ref:02d}", freq="M")
            parte = mensual[
                (mensual["POZO_KEY"] == pozo_key) &
                (mensual["ANIO_MES"] == periodo_ref)
            ]

            if parte.empty:
                prcr_mes = 0.0
                prag_mes = 0.0
                interv_mes = 0
            else:
                prcr_mes = float(parte["PRCR_MES"].iloc[0])
                prag_mes = float(parte["PRAG_MES"].iloc[0])
                interv_mes = int(parte["INTERV_MES"].iloc[0])

            dias = dias_mes(a_ref, m_ref)
            potencial = prcr_mes / dias if dias > 0 else 0

            if offset == 0:
                registro["DIAS_ULTIMO_MES_INTERVENIDO"] = dias
                registro["PRCR_ULTIMO_MES_INTERVENIDO"] = prcr_mes
                registro["PRAG_ULTIMO_MES_INTERVENIDO"] = prag_mes
                registro["INTERV_ULTIMO_MES_INTERVENIDO"] = interv_mes
                registro["POTENCIAL_ULTIMO_MES_INTERVENIDO_BOPD"] = potencial
            else:
                registro[f"MES_{etiqueta}"] = periodo_texto(a_ref, m_ref)
                registro[f"DIAS_{etiqueta}"] = dias
                registro[f"PRCR_{etiqueta}"] = prcr_mes
                registro[f"PRAG_{etiqueta}"] = prag_mes
                registro[f"INTERV_{etiqueta}"] = interv_mes
                registro[f"POTENCIAL_{etiqueta}_BOPD"] = potencial

        registro["POTENCIAL_PROMEDIO_3_MESES_INTERVENCION_BOPD"] = (
            registro["POTENCIAL_ULTIMO_MES_INTERVENIDO_BOPD"] +
            registro["POTENCIAL_1_MES_ANTES_ULTIMA_INTERVENCION_BOPD"] +
            registro["POTENCIAL_2_MES_ANTES_ULTIMA_INTERVENCION_BOPD"]
        ) / 3

        registros.append(registro)

    potencial_interv = pd.DataFrame(registros)

    # Última fecha con producción real de petróleo y la fecha productiva anterior.
    prod_eventos = data[data["PRCR"] > 0].sort_values(["POZO_KEY", "FECHA"]).copy()

    if not prod_eventos.empty:
        ult_prod = (
            prod_eventos.groupby("POZO_KEY", as_index=False)
            .tail(1)[["POZO_KEY", "FECHA", "PRCR", "PRAG"]]
            .rename(columns={
                "FECHA": "ULTIMA_FECHA_CON_PRODUCCION_PETROLEO",
                "PRCR": "PRODUCCION_PETROLEO_ULTIMA_FECHA",
                "PRAG": "PRODUCCION_AGUA_ULTIMA_FECHA"
            })
        )

        anteriores = []
        for pozo_key, grupo in prod_eventos.groupby("POZO_KEY"):
            if len(grupo) >= 2:
                fila = grupo.iloc[-2]
                anteriores.append({
                    "POZO_KEY": pozo_key,
                    "FECHA_PRODUCCION_PETROLEO_ANTERIOR": fila["FECHA"],
                    "PRODUCCION_PETROLEO_FECHA_ANTERIOR": float(fila["PRCR"]),
                    "PRODUCCION_AGUA_FECHA_ANTERIOR": float(fila["PRAG"])
                })

        prod_ant = pd.DataFrame(anteriores)

        if prod_ant.empty:
            prod_ant = pd.DataFrame(columns=[
                "POZO_KEY",
                "FECHA_PRODUCCION_PETROLEO_ANTERIOR",
                "PRODUCCION_PETROLEO_FECHA_ANTERIOR",
                "PRODUCCION_AGUA_FECHA_ANTERIOR"
            ])

        ult_prod["ANIO_MES_PROD"] = ult_prod["ULTIMA_FECHA_CON_PRODUCCION_PETROLEO"].dt.to_period("M")

        mensual_prod = mensual.rename(columns={
            "ANIO_MES": "ANIO_MES_PROD",
            "PRCR_MES": "PRCR_ULTIMO_MES_CON_PRODUCCION",
            "PRAG_MES": "PRAG_ULTIMO_MES_CON_PRODUCCION",
            "INTERV_MES": "INTERV_ULTIMO_MES_CON_PRODUCCION"
        })[[
            "POZO_KEY",
            "ANIO_MES_PROD",
            "PRCR_ULTIMO_MES_CON_PRODUCCION",
            "PRAG_ULTIMO_MES_CON_PRODUCCION",
            "INTERV_ULTIMO_MES_CON_PRODUCCION"
        ]]

        ult_prod = ult_prod.merge(
            mensual_prod,
            on=["POZO_KEY", "ANIO_MES_PROD"],
            how="left"
        )

        ult_prod["ANIO_PROD"] = ult_prod["ANIO_MES_PROD"].dt.year
        ult_prod["MES_PROD"] = ult_prod["ANIO_MES_PROD"].dt.month
        ult_prod["ULTIMO_MES_CON_PRODUCCION_PETROLEO"] = ult_prod.apply(
            lambda r: periodo_texto(r["ANIO_PROD"], r["MES_PROD"]),
            axis=1
        )
        ult_prod["DIAS_ULTIMO_MES_CON_PRODUCCION"] = ult_prod.apply(
            lambda r: dias_mes(r["ANIO_PROD"], r["MES_PROD"]),
            axis=1
        )
        ult_prod["POTENCIAL_ULTIMO_MES_CON_PRODUCCION_BOPD"] = np.where(
            ult_prod["DIAS_ULTIMO_MES_CON_PRODUCCION"] > 0,
            ult_prod["PRCR_ULTIMO_MES_CON_PRODUCCION"] / ult_prod["DIAS_ULTIMO_MES_CON_PRODUCCION"],
            0
        )

        prod_trazabilidad = ult_prod.merge(prod_ant, on="POZO_KEY", how="left")

        prod_trazabilidad = prod_trazabilidad[[
            "POZO_KEY",
            "ULTIMA_FECHA_CON_PRODUCCION_PETROLEO",
            "PRODUCCION_PETROLEO_ULTIMA_FECHA",
            "PRODUCCION_AGUA_ULTIMA_FECHA",
            "FECHA_PRODUCCION_PETROLEO_ANTERIOR",
            "PRODUCCION_PETROLEO_FECHA_ANTERIOR",
            "PRODUCCION_AGUA_FECHA_ANTERIOR",
            "ULTIMO_MES_CON_PRODUCCION_PETROLEO",
            "DIAS_ULTIMO_MES_CON_PRODUCCION",
            "PRCR_ULTIMO_MES_CON_PRODUCCION",
            "PRAG_ULTIMO_MES_CON_PRODUCCION",
            "INTERV_ULTIMO_MES_CON_PRODUCCION",
            "POTENCIAL_ULTIMO_MES_CON_PRODUCCION_BOPD"
        ]]
    else:
        prod_trazabilidad = pd.DataFrame(columns=[
            "POZO_KEY",
            "ULTIMA_FECHA_CON_PRODUCCION_PETROLEO",
            "PRODUCCION_PETROLEO_ULTIMA_FECHA",
            "PRODUCCION_AGUA_ULTIMA_FECHA",
            "FECHA_PRODUCCION_PETROLEO_ANTERIOR",
            "PRODUCCION_PETROLEO_FECHA_ANTERIOR",
            "PRODUCCION_AGUA_FECHA_ANTERIOR",
            "ULTIMO_MES_CON_PRODUCCION_PETROLEO",
            "DIAS_ULTIMO_MES_CON_PRODUCCION",
            "PRCR_ULTIMO_MES_CON_PRODUCCION",
            "PRAG_ULTIMO_MES_CON_PRODUCCION",
            "INTERV_ULTIMO_MES_CON_PRODUCCION",
            "POTENCIAL_ULTIMO_MES_CON_PRODUCCION_BOPD"
        ])

    out = universo.merge(potencial_interv, on="POZO_KEY", how="left")
    out = out.merge(prod_trazabilidad, on="POZO_KEY", how="left")

    texto_cols = [
        "ULTIMO_MES_INTERVENIDO",
        "MES_1_MES_ANTES_ULTIMA_INTERVENCION",
        "MES_2_MES_ANTES_ULTIMA_INTERVENCION",
        "ULTIMO_MES_CON_PRODUCCION_PETROLEO"
    ]

    for col in texto_cols:
        if col in out.columns:
            out[col] = out[col].fillna("Sin dato")

    fecha_cols = [
        "ULTIMA_FECHA_INTERVENCION",
        "ULTIMA_FECHA_CON_PRODUCCION_PETROLEO",
        "FECHA_PRODUCCION_PETROLEO_ANTERIOR"
    ]

    for col in fecha_cols:
        if col in out.columns:
            out[col] = pd.to_datetime(out[col], errors="coerce")

    num_cols = [
        "DIAS_ULTIMO_MES_INTERVENIDO",
        "PRCR_ULTIMO_MES_INTERVENIDO",
        "PRAG_ULTIMO_MES_INTERVENIDO",
        "INTERV_ULTIMO_MES_INTERVENIDO",
        "POTENCIAL_ULTIMO_MES_INTERVENIDO_BOPD",
        "DIAS_1_MES_ANTES_ULTIMA_INTERVENCION",
        "PRCR_1_MES_ANTES_ULTIMA_INTERVENCION",
        "PRAG_1_MES_ANTES_ULTIMA_INTERVENCION",
        "INTERV_1_MES_ANTES_ULTIMA_INTERVENCION",
        "POTENCIAL_1_MES_ANTES_ULTIMA_INTERVENCION_BOPD",
        "DIAS_2_MES_ANTES_ULTIMA_INTERVENCION",
        "PRCR_2_MES_ANTES_ULTIMA_INTERVENCION",
        "PRAG_2_MES_ANTES_ULTIMA_INTERVENCION",
        "INTERV_2_MES_ANTES_ULTIMA_INTERVENCION",
        "POTENCIAL_2_MES_ANTES_ULTIMA_INTERVENCION_BOPD",
        "POTENCIAL_PROMEDIO_3_MESES_INTERVENCION_BOPD",
        "PRODUCCION_PETROLEO_ULTIMA_FECHA",
        "PRODUCCION_AGUA_ULTIMA_FECHA",
        "PRODUCCION_PETROLEO_FECHA_ANTERIOR",
        "PRODUCCION_AGUA_FECHA_ANTERIOR",
        "DIAS_ULTIMO_MES_CON_PRODUCCION",
        "PRCR_ULTIMO_MES_CON_PRODUCCION",
        "PRAG_ULTIMO_MES_CON_PRODUCCION",
        "INTERV_ULTIMO_MES_CON_PRODUCCION",
        "POTENCIAL_ULTIMO_MES_CON_PRODUCCION_BOPD"
    ]

    for col in num_cols:
        if col not in out.columns:
            out[col] = 0
        out[col] = out[col].fillna(0)

    return out.sort_values("POTENCIAL_PROMEDIO_3_MESES_INTERVENCION_BOPD", ascending=False)

def calcular_analisis(df, universo, anio, mes, meses_base, baterias, tipos, clases_general, clases_desplazadas):
    # Data del mes objetivo
    mes_df = df[(df["ANIO"] == int(anio)) & (df["MES"] == int(mes))].copy()

    # Filtros generales
    if baterias:
        mes_df = mes_df[mes_df["BATERIA"].isin(baterias)]
    if tipos:
        mes_df = mes_df[mes_df["TIPO_SWAB"].isin(tipos)]
    if clases_general:
        mes_df = mes_df[mes_df["CLASIFICACION"].isin(clases_general)]

    universo_general = universo.copy()
    if baterias:
        universo_general = universo_general[universo_general["BATERIA"].isin(baterias)]
    if clases_general:
        universo_general = universo_general[universo_general["CLASIFICACION"].isin(clases_general)]

    resumen_mes = (
        mes_df.groupby("POZO_KEY", as_index=False)
        .agg(
            INTERVENCIONES=("FECHA", "count"),
            PRCR=("PRCR", "sum"),
            PRAG=("PRAG", "sum"),
            TIPO_SWAB=("TIPO_SWAB", lambda x: ", ".join(sorted(set(x)))),
            ULTIMA_FECHA=("FECHA", "max")
        )
    )

    resumen_pozos = universo_general.merge(resumen_mes, on="POZO_KEY", how="left")
    resumen_pozos["INTERVENCIONES"] = resumen_pozos["INTERVENCIONES"].fillna(0).astype(int)
    resumen_pozos["PRCR"] = resumen_pozos["PRCR"].fillna(0)
    resumen_pozos["PRAG"] = resumen_pozos["PRAG"].fillna(0)
    resumen_pozos["TIPO_SWAB"] = resumen_pozos["TIPO_SWAB"].fillna("")
    resumen_pozos["ESTADO"] = np.where(resumen_pozos["INTERVENCIONES"] > 0, "Intervenido", "No intervenido")
    resumen_pozos["OIL_POR_INTERV"] = np.where(resumen_pozos["INTERVENCIONES"] > 0, resumen_pozos["PRCR"] / resumen_pozos["INTERVENCIONES"], 0)
    resumen_pozos["AGUA_POR_INTERV"] = np.where(resumen_pozos["INTERVENCIONES"] > 0, resumen_pozos["PRAG"] / resumen_pozos["INTERVENCIONES"], 0)

    # Potencial
    potencial = calcular_potencial(df, universo, baterias, tipos, clases_general if clases_general else None)

    # Base N meses
    meses_previos = lista_meses_previos(anio, mes, meses_base)
    partes_base = []
    for a, m in meses_previos:
        partes_base.append(df[(df["ANIO"] == a) & (df["MES"] == m)].copy())
    base = pd.concat(partes_base, ignore_index=True) if partes_base else df.iloc[0:0].copy()

    # Mes anterior directo
    a_ant, m_ant = mes_anterior(anio, mes, 1)
    mes_ant = df[(df["ANIO"] == a_ant) & (df["MES"] == m_ant)].copy()

    if baterias:
        base = base[base["BATERIA"].isin(baterias)]
        mes_ant = mes_ant[mes_ant["BATERIA"].isin(baterias)]
    if tipos:
        base = base[base["TIPO_SWAB"].isin(tipos)]
        mes_ant = mes_ant[mes_ant["TIPO_SWAB"].isin(tipos)]

    universo_desp = universo.copy()
    if baterias:
        universo_desp = universo_desp[universo_desp["BATERIA"].isin(baterias)]
    if clases_desplazadas:
        universo_desp = universo_desp[universo_desp["CLASIFICACION"].isin(clases_desplazadas)]

    actual_desp = df[(df["ANIO"] == int(anio)) & (df["MES"] == int(mes))].copy()
    if baterias:
        actual_desp = actual_desp[actual_desp["BATERIA"].isin(baterias)]
    if tipos:
        actual_desp = actual_desp[actual_desp["TIPO_SWAB"].isin(tipos)]

    res_base = resumen_pozo_periodo(base, "BASE", divisor=meses_base)
    res_actual = resumen_pozo_periodo(actual_desp, "ACTUAL", divisor=1)
    res_ant = resumen_pozo_periodo(mes_ant, "MES_ANTERIOR_OBJETIVO", divisor=1)

    desp = universo_desp.merge(res_base, on="POZO_KEY", how="left")
    desp = desp.merge(res_actual, on="POZO_KEY", how="left")
    desp = desp.merge(res_ant, on="POZO_KEY", how="left")

    for c in [
        "INTERV_BASE", "PRCR_BASE", "PRAG_BASE", "OIL_INTERV_BASE",
        "INTERV_ACTUAL", "PRCR_ACTUAL", "PRAG_ACTUAL", "OIL_INTERV_ACTUAL",
        "INTERV_MES_ANTERIOR_OBJETIVO", "PRCR_MES_ANTERIOR_OBJETIVO", "PRAG_MES_ANTERIOR_OBJETIVO"
    ]:
        if c not in desp.columns:
            desp[c] = 0
        desp[c] = desp[c].fillna(0)

    desp["VAR_INTERV"] = desp["INTERV_ACTUAL"] - desp["INTERV_BASE"]
    desp["VAR_PRCR"] = desp["PRCR_ACTUAL"] - desp["PRCR_BASE"]
    desp["VAR_PRAG"] = desp["PRAG_ACTUAL"] - desp["PRAG_BASE"]

    desp["INTERV_BASE_NO_REALIZADAS"] = np.where(desp["VAR_INTERV"] < 0, -desp["VAR_INTERV"], 0)
    desp["PRCR_BASE_NO_REALIZADO"] = np.where(desp["VAR_PRCR"] < 0, -desp["VAR_PRCR"], 0)
    desp["PRAG_BASE_NO_REALIZADO"] = np.where(desp["VAR_PRAG"] < 0, -desp["VAR_PRAG"], 0)

    desp["VAR_PETROLEO_VS_MES_ANTERIOR"] = desp["PRCR_ACTUAL"] - desp["PRCR_MES_ANTERIOR_OBJETIVO"]
    desp["PRODUCCION_PETROLEO_DEJADA_MES_ANTERIOR"] = np.where(
        desp["VAR_PETROLEO_VS_MES_ANTERIOR"] < 0,
        -desp["VAR_PETROLEO_VS_MES_ANTERIOR"],
        0
    )

    def estado(r):
        if r["INTERV_BASE"] > 0 and r["INTERV_ACTUAL"] == 0:
            return "DEJADO DE HACER"
        if r["INTERV_BASE"] > 0 and r["INTERV_ACTUAL"] > 0 and r["INTERV_ACTUAL"] < r["INTERV_BASE"]:
            return "REDUCIDO"
        if r["INTERV_BASE"] == 0 and r["INTERV_ACTUAL"] > 0:
            return "NUEVO / RETOMADO"
        if r["INTERV_BASE"] > 0 and r["INTERV_ACTUAL"] >= r["INTERV_BASE"]:
            return "MANTENIDO / AUMENTADO"
        return "SIN ACTIVIDAD"

    desp["ESTADO_DESPLAZAMIENTO"] = desp.apply(estado, axis=1)

    candidatos = desp[desp["ESTADO_DESPLAZAMIENTO"].isin(["DEJADO DE HACER", "REDUCIDO"])]
    q75 = candidatos["PRCR_BASE_NO_REALIZADO"].quantile(0.75) if not candidatos.empty else 0
    q50 = candidatos["PRCR_BASE_NO_REALIZADO"].quantile(0.50) if not candidatos.empty else 0

    desp["PRIORIDAD_REVISION"] = np.where(
        (desp["PRCR_BASE_NO_REALIZADO"] >= q75) & (desp["PRCR_BASE_NO_REALIZADO"] > 0),
        "ALTA",
        np.where(
            (desp["PRCR_BASE_NO_REALIZADO"] >= q50) & (desp["PRCR_BASE_NO_REALIZADO"] > 0),
            "MEDIA",
            np.where(desp["ESTADO_DESPLAZAMIENTO"].isin(["DEJADO DE HACER", "REDUCIDO"]), "BAJA", "")
        )
    )

    pot_cols = [
        "POZO_KEY",
        "ULTIMO_MES_INTERVENIDO",
        "ULTIMA_FECHA_INTERVENCION",
        "DIAS_ULTIMO_MES_INTERVENIDO",
        "PRCR_ULTIMO_MES_INTERVENIDO",
        "PRAG_ULTIMO_MES_INTERVENIDO",
        "INTERV_ULTIMO_MES_INTERVENIDO",
        "POTENCIAL_ULTIMO_MES_INTERVENIDO_BOPD",
        "MES_1_MES_ANTES_ULTIMA_INTERVENCION",
        "DIAS_1_MES_ANTES_ULTIMA_INTERVENCION",
        "PRCR_1_MES_ANTES_ULTIMA_INTERVENCION",
        "PRAG_1_MES_ANTES_ULTIMA_INTERVENCION",
        "INTERV_1_MES_ANTES_ULTIMA_INTERVENCION",
        "POTENCIAL_1_MES_ANTES_ULTIMA_INTERVENCION_BOPD",
        "MES_2_MES_ANTES_ULTIMA_INTERVENCION",
        "DIAS_2_MES_ANTES_ULTIMA_INTERVENCION",
        "PRCR_2_MES_ANTES_ULTIMA_INTERVENCION",
        "PRAG_2_MES_ANTES_ULTIMA_INTERVENCION",
        "INTERV_2_MES_ANTES_ULTIMA_INTERVENCION",
        "POTENCIAL_2_MES_ANTES_ULTIMA_INTERVENCION_BOPD",
        "POTENCIAL_PROMEDIO_3_MESES_INTERVENCION_BOPD",
        "ULTIMA_FECHA_CON_PRODUCCION_PETROLEO",
        "PRODUCCION_PETROLEO_ULTIMA_FECHA",
        "PRODUCCION_AGUA_ULTIMA_FECHA",
        "FECHA_PRODUCCION_PETROLEO_ANTERIOR",
        "PRODUCCION_PETROLEO_FECHA_ANTERIOR",
        "PRODUCCION_AGUA_FECHA_ANTERIOR",
        "ULTIMO_MES_CON_PRODUCCION_PETROLEO",
        "DIAS_ULTIMO_MES_CON_PRODUCCION",
        "PRCR_ULTIMO_MES_CON_PRODUCCION",
        "PRAG_ULTIMO_MES_CON_PRODUCCION",
        "INTERV_ULTIMO_MES_CON_PRODUCCION",
        "POTENCIAL_ULTIMO_MES_CON_PRODUCCION_BOPD"
    ]
    desp = desp.merge(potencial[pot_cols], on="POZO_KEY", how="left")

    orden = {"DEJADO DE HACER": 1, "REDUCIDO": 2, "NUEVO / RETOMADO": 3, "MANTENIDO / AUMENTADO": 4, "SIN ACTIVIDAD": 5}
    desp["ORDEN"] = desp["ESTADO_DESPLAZAMIENTO"].map(orden).fillna(9)
    desp = desp.sort_values(["ORDEN", "PRCR_BASE_NO_REALIZADO"], ascending=[True, False]).drop(columns="ORDEN")

    res_clase = resumen_pozos.groupby("CLASIFICACION", as_index=False).agg(
        POZOS_TOTAL=("POZO_KEY", "nunique"),
        POZOS_INTERVENIDOS=("ESTADO", lambda x: (x == "Intervenido").sum()),
        POZOS_NO_INTERVENIDOS=("ESTADO", lambda x: (x == "No intervenido").sum()),
        INTERVENCIONES=("INTERVENCIONES", "sum"),
        PRCR=("PRCR", "sum"),
        PRAG=("PRAG", "sum")
    )

    res_bat = resumen_pozos.groupby("BATERIA", as_index=False).agg(
        POZOS_TOTAL=("POZO_KEY", "nunique"),
        POZOS_INTERVENIDOS=("ESTADO", lambda x: (x == "Intervenido").sum()),
        POZOS_NO_INTERVENIDOS=("ESTADO", lambda x: (x == "No intervenido").sum()),
        INTERVENCIONES=("INTERVENCIONES", "sum"),
        PRCR=("PRCR", "sum"),
        PRAG=("PRAG", "sum")
    ).sort_values("PRCR", ascending=False)

    if mes_df.empty:
        res_tipo = pd.DataFrame(columns=["TIPO_SWAB", "POZOS", "INTERVENCIONES", "PRCR", "PRAG"])
    else:
        res_tipo = mes_df.groupby("TIPO_SWAB", as_index=False).agg(
            POZOS=("POZO_KEY", "nunique"),
            INTERVENCIONES=("FECHA", "count"),
            PRCR=("PRCR", "sum"),
            PRAG=("PRAG", "sum")
        )

    tendencia = df[df["ANIO"] == int(anio)].groupby(["MES", "MES_NOMBRE"], as_index=False).agg(
        POZOS=("POZO_KEY", "nunique"),
        INTERVENCIONES=("FECHA", "count"),
        PRCR=("PRCR", "sum"),
        PRAG=("PRAG", "sum")
    ).sort_values("MES")

    list_conv = universo[universo["CLASIFICACION"].str.startswith("Convertido")].merge(
        resumen_pozos[["POZO_KEY", "ESTADO", "INTERVENCIONES", "PRCR", "PRAG", "TIPO_SWAB"]],
        on="POZO_KEY",
        how="left"
    )
    list_conv["ESTADO"] = list_conv["ESTADO"].fillna("No intervenido")
    for c in ["INTERVENCIONES", "PRCR", "PRAG"]:
        list_conv[c] = list_conv[c].fillna(0)

    return {
        "resumen_pozos": resumen_pozos,
        "desp": desp,
        "potencial": potencial,
        "res_clase": res_clase,
        "res_bat": res_bat,
        "res_tipo": res_tipo,
        "tendencia": tendencia,
        "list_conv": list_conv,
        "actual": actual_desp,
        "base": base,
        "meses_previos": meses_previos
    }


# ============================================================
# INTERFAZ
# ============================================================

st.title("🛢️ SWAB Lote X - Servicio y Potencial")
st.caption("Análisis rápido de pozos dejados, impacto de convertidos 2026, potencial mensual y exclusión de candidatos ATA.")

archivo = st.file_uploader("Sube el Excel principal con la hoja Datos de Swab", type=["xlsx"])

if archivo is None:
    st.info("Sube el Excel para iniciar.")
    st.stop()

with st.spinner("Cargando Excel. Esta primera carga puede demorar unos segundos..."):
    try:
        df_total, df_ata, hoja_usada = cargar_datos(archivo.getvalue())
    except Exception as e:
        st.error(f"No se pudo cargar el Excel: {e}")
        st.stop()

fecha_min = df_total["FECHA"].min().date()
fecha_max = df_total["FECHA"].max().date()

with st.sidebar:
    st.header("1. Rango de análisis")
    fecha_inicio = st.date_input("Desde", value=fecha_min, min_value=fecha_min, max_value=fecha_max)
    fecha_fin = st.date_input("Hasta", value=fecha_max, min_value=fecha_min, max_value=fecha_max)

df = df_total[(df_total["FECHA"] >= pd.to_datetime(fecha_inicio)) & (df_total["FECHA"] <= pd.to_datetime(fecha_fin))].copy()

if df.empty:
    st.error("El rango seleccionado no tiene datos.")
    st.stop()

universo = construir_universo(df)
anios = sorted(df["ANIO"].unique().astype(int).tolist())
anio_def, mes_def = ultimo_mes_completo(df)
baterias = sorted([x for x in universo["BATERIA"].dropna().unique().tolist() if x != ""])
tipos = sorted([x for x in df["TIPO_SWAB"].dropna().unique().tolist() if x != ""])
clases = ["Básica", "Convertido 2024", "Convertido 2025", "Convertido 2026"]

with st.sidebar:
    st.header("2. Módulo principal")
    modulo = st.radio(
        "Selecciona con qué trabajar",
        ["Análisis de pozos dejados", "Potencial mensual de pozos", "Vista completa"],
        index=2
    )

    st.header("3. Periodo objetivo")
    anio_obj = st.selectbox(
        "Año objetivo",
        anios,
        index=anios.index(anio_def) if anio_def in anios else len(anios) - 1
    )

    meses_disp = sorted(df[df["ANIO"] == anio_obj]["MES"].unique().astype(int).tolist())
    mes_obj = st.selectbox(
        "Mes objetivo",
        meses_disp,
        index=meses_disp.index(mes_def) if mes_def in meses_disp else len(meses_disp) - 1,
        format_func=lambda x: MESES[int(x)]
    )

    meses_base = st.slider("Comparar contra N meses anteriores", 1, 15, 3, 1)

    st.header("4. Filtros operativos")
    baterias_sel = st.multiselect("Batería", baterias, default=[])
    tipos_sel = st.multiselect("Tipo de swab", tipos, default=[])
    clases_sel = st.multiselect("Tipo de pozo para vistas generales", clases, default=[])
    clases_desplazadas = st.multiselect(
        "Pozos a evaluar como desplazados",
        clases,
        default=["Básica", "Convertido 2024", "Convertido 2025"]
    )
    top_n = st.slider("Top para gráficos", 5, 50, 20, 5)

    ejecutar = st.button("Ejecutar análisis", type="primary")

if not ejecutar and "res_fast" not in st.session_state:
    st.warning("Configura los filtros y presiona Ejecutar análisis.")
    st.stop()

if ejecutar:
    with st.spinner("Ejecutando análisis. Espere un momento..."):
        st.session_state["res_fast"] = calcular_analisis(
            df=df,
            universo=universo,
            anio=anio_obj,
            mes=mes_obj,
            meses_base=meses_base,
            baterias=baterias_sel,
            tipos=tipos_sel,
            clases_general=clases_sel,
            clases_desplazadas=clases_desplazadas
        )
        st.session_state["param_fast"] = {
            "anio": anio_obj,
            "mes": mes_obj,
            "meses_base": meses_base,
            "top_n": top_n,
            "fecha_inicio": fecha_inicio,
            "fecha_fin": fecha_fin
        }

res = st.session_state["res_fast"]
params = st.session_state["param_fast"]

anio_obj = params["anio"]
mes_obj = params["mes"]
meses_base = params["meses_base"]
top_n = params["top_n"]

res_pozos = res["resumen_pozos"]
desp = res["desp"]
pot = res["potencial"]
res_clase = res["res_clase"]
res_bat = res["res_bat"]
res_tipo = res["res_tipo"]
res_tend = res["tendencia"]
list_conv = res["list_conv"]
actual = res["actual"]
meses_previos = res["meses_previos"]

periodo = periodo_texto(anio_obj, mes_obj)

if fecha_max < ultimo_dia_mes(anio_obj, mes_obj).date() and anio_obj == fecha_max.year and mes_obj == fecha_max.month:
    st.warning(f"El mes {periodo} está incompleto. La data llega hasta {fecha_max}.")

st.subheader(f"Resumen ejecutivo: {periodo}")

pozos_universo = res_pozos["POZO_KEY"].nunique()
pozos_interv = int((res_pozos["ESTADO"] == "Intervenido").sum())
pozos_no = int((res_pozos["ESTADO"] == "No intervenido").sum())
interv = int(res_pozos["INTERVENCIONES"].sum())
petroleo = float(res_pozos["PRCR"].sum())
agua = float(res_pozos["PRAG"].sum())

c1, c2, c3, c4, c5, c6 = st.columns(6)
c1.metric("Pozos universo", f"{pozos_universo:,}")
c2.metric("Intervenidos", f"{pozos_interv:,}")
c3.metric("No intervenidos", f"{pozos_no:,}")
c4.metric("Intervenciones", f"{interv:,}")
c5.metric("Producción de petróleo", f"{petroleo:,.2f}")
c6.metric("Producción de agua", f"{agua:,.2f}")

st.caption(
    f"Hoja usada: {hoja_usada}. Rango analizado: {params['fecha_inicio']} al {params['fecha_fin']}. "
    f"Pozos ATA excluidos encontrados en la data: {df_ata['POZO_KEY'].nunique()}."
)

st.divider()

tabs = st.tabs([
    "Análisis de pozos dejados",
    "Impacto convertidos 2026",
    "Listado convertidos",
    "Potencial mensual",
    "Pozos y baterías",
    "TS y CS",
    "Descargas"
])

with tabs[0]:
    st.subheader("Análisis de pozos dejados de hacer o reducidos")

    dejados = desp[desp["ESTADO_DESPLAZAMIENTO"] == "DEJADO DE HACER"]
    reducidos = desp[desp["ESTADO_DESPLAZAMIENTO"] == "REDUCIDO"]

    k1, k2, k3, k4, k5, k6 = st.columns(6)
    k1.metric("Pozos dejados", f"{len(dejados):,}")
    k2.metric("Pozos reducidos", f"{len(reducidos):,}")
    k3.metric("Interv. base no realizadas", f"{desp['INTERV_BASE_NO_REALIZADAS'].sum():,.2f}")
    k4.metric("Petróleo base no realizado", f"{desp['PRCR_BASE_NO_REALIZADO'].sum():,.2f}")
    k5.metric("Petróleo dejado vs mes anterior", f"{desp['PRODUCCION_PETROLEO_DEJADA_MES_ANTERIOR'].sum():,.2f}")
    k6.metric("Potencial promedio 3 meses", f"{desp['POTENCIAL_PROMEDIO_3_MESES_INTERVENCION_BOPD'].sum():,.2f} bopd")

    base_txt = ", ".join([periodo_texto(a, m) for a, m in meses_previos])
    st.info(f"Periodo objetivo: {periodo}. Base comparativa: promedio mensual de {meses_base} mes(es) anteriores: {base_txt}.")

    estado = st.radio(
        "Estado de desplazamiento",
        ["DEJADO DE HACER", "REDUCIDO", "NUEVO / RETOMADO", "MANTENIDO / AUMENTADO", "SIN ACTIVIDAD", "TODOS"],
        horizontal=True
    )

    tabla = desp.copy()
    if estado != "TODOS":
        tabla = tabla[tabla["ESTADO_DESPLAZAMIENTO"] == estado]

    cols_desp = [
        "ESTADO_DESPLAZAMIENTO", "PRIORIDAD_REVISION", "POZO", "BATERIA", "CLASIFICACION", "ANIO_CONVERSION",
        "INTERV_BASE", "INTERV_ACTUAL", "INTERV_BASE_NO_REALIZADAS",
        "PRCR_BASE", "PRCR_ACTUAL", "PRCR_BASE_NO_REALIZADO",
        "PRCR_MES_ANTERIOR_OBJETIVO", "PRODUCCION_PETROLEO_DEJADA_MES_ANTERIOR", "VAR_PETROLEO_VS_MES_ANTERIOR",
        "ULTIMO_MES_INTERVENIDO", "ULTIMA_FECHA_INTERVENCION",
        "PRCR_ULTIMO_MES_INTERVENIDO", "POTENCIAL_ULTIMO_MES_INTERVENIDO_BOPD",
        "MES_1_MES_ANTES_ULTIMA_INTERVENCION",
        "PRCR_1_MES_ANTES_ULTIMA_INTERVENCION",
        "POTENCIAL_1_MES_ANTES_ULTIMA_INTERVENCION_BOPD",
        "MES_2_MES_ANTES_ULTIMA_INTERVENCION",
        "PRCR_2_MES_ANTES_ULTIMA_INTERVENCION",
        "POTENCIAL_2_MES_ANTES_ULTIMA_INTERVENCION_BOPD",
        "POTENCIAL_PROMEDIO_3_MESES_INTERVENCION_BOPD",
        "ULTIMA_FECHA_CON_PRODUCCION_PETROLEO",
        "PRODUCCION_PETROLEO_ULTIMA_FECHA",
        "FECHA_PRODUCCION_PETROLEO_ANTERIOR",
        "PRODUCCION_PETROLEO_FECHA_ANTERIOR",
        "ULTIMO_MES_CON_PRODUCCION_PETROLEO",
        "PRCR_ULTIMO_MES_CON_PRODUCCION",
        "POTENCIAL_ULTIMO_MES_CON_PRODUCCION_BOPD",
        "PRAG_BASE", "PRAG_ACTUAL", "PRAG_BASE_NO_REALIZADO"
    ]
    cols_desp = [c for c in cols_desp if c in tabla.columns]

    st.dataframe(vista_tabla(tabla[cols_desp]), use_container_width=True, hide_index=True)

    top_dej = (
        desp[desp["ESTADO_DESPLAZAMIENTO"].isin(["DEJADO DE HACER", "REDUCIDO"])]
        .sort_values("PRCR_BASE_NO_REALIZADO", ascending=False)
        .head(top_n)
    )

    if not top_dej.empty:
        fig = px.bar(
            top_dej,
            x="POZO",
            y="PRCR_BASE_NO_REALIZADO",
            color="ESTADO_DESPLAZAMIENTO",
            text="PRCR_BASE_NO_REALIZADO",
            hover_data=["BATERIA", "CLASIFICACION", "POTENCIAL_PROMEDIO_3_MESES_INTERVENCION_BOPD"]
        )
        st.plotly_chart(aplicar_layout(fig, f"Top {top_n} pozos con petróleo base no realizado", 560), use_container_width=True)

        fig2 = px.scatter(
            top_dej,
            x="POTENCIAL_PROMEDIO_3_MESES_INTERVENCION_BOPD",
            y="PRCR_BASE_NO_REALIZADO",
            size="INTERV_BASE_NO_REALIZADAS",
            color="CLASIFICACION",
            hover_name="POZO",
            hover_data=["BATERIA", "ESTADO_DESPLAZAMIENTO", "PRODUCCION_PETROLEO_DEJADA_MES_ANTERIOR"]
        )
        st.plotly_chart(aplicar_layout(fig2, "Potencial promedio 3 meses vs petróleo base no realizado", 560), use_container_width=True)

    st.download_button(
        "Descargar data de gráficas de pozos dejados",
        data=excel_descarga({"Top dejados reducidos": top_dej, "Tabla desplazamiento": desp[cols_desp]}),
        file_name=f"data_graficas_pozos_dejados_{anio_obj}_{mes_obj}.xlsx",
        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    )

with tabs[1]:
    st.subheader("Impacto de priorizar convertidos 2026")

    conv_actual = actual[actual["CLASIFICACION"] == "Convertido 2026"]
    conv_det = _conv2026_detalle(list_conv, actual, pot)

    c1, c2, c3, c4, c5 = st.columns(5)
    c1.metric("Conv. 2026 intervenidos", f"{conv_actual['POZO_KEY'].nunique():,} de 20")
    c2.metric("Interv. conv. 2026", f"{len(conv_actual):,}")
    c3.metric("Producción de petróleo", f"{conv_actual['PRCR'].sum():,.2f}")
    c4.metric("Producción de agua", f"{conv_actual['PRAG'].sum():,.2f}")
    c5.metric(
        "Petróleo por intervención",
        f"{(conv_actual['PRCR'].sum() / len(conv_actual)):,.2f}" if len(conv_actual) > 0 else "0.00"
    )

    st.caption(
        "Esta vista separa los 20 convertidos 2026 como grupo priorizado y muestra su producción, "
        "frecuencia, potencial y trazabilidad de la última producción real de petróleo."
    )

    if not conv_det.empty:
        col_a, col_b = st.columns(2)

        with col_a:
            fig_conv_pet = px.bar(
                conv_det.sort_values("PRCR", ascending=False).head(top_n),
                x="POZO",
                y="PRCR",
                color="ESTADO_CONVERTIDO_2026",
                text="PRCR",
                hover_data=["BATERIA", "INTERVENCIONES", "PRAG", "PETROLEO_POR_INTERVENCION"]
            )
            st.plotly_chart(
                aplicar_layout(fig_conv_pet, f"Top {top_n} convertidos 2026 por producción de petróleo", 520),
                use_container_width=True
            )

        with col_b:
            fig_conv_pot = px.scatter(
                conv_det,
                x="POTENCIAL_PROMEDIO_3_MESES_INTERVENCION_BOPD",
                y="PRCR",
                size="INTERVENCIONES",
                color="ESTADO_CONVERTIDO_2026",
                hover_name="POZO",
                hover_data=[
                    "BATERIA",
                    "ULTIMO_MES_INTERVENIDO",
                    "ULTIMA_FECHA_CON_PRODUCCION_PETROLEO",
                    "PRODUCCION_PETROLEO_ULTIMA_FECHA"
                ]
            )
            st.plotly_chart(
                aplicar_layout(fig_conv_pot, "Convertidos 2026: potencial vs producción actual", 520),
                use_container_width=True
            )

        fig_estado = px.pie(
            conv_det,
            names="ESTADO_CONVERTIDO_2026",
            values="INTERVENCIONES",
            hole=0.45,
            title="Distribución de intervenciones en convertidos 2026"
        )
        fig_estado.update_traces(texttemplate="%{label}<br>%{percent:.2%}")
        st.plotly_chart(aplicar_layout(fig_estado, "Distribución de intervenciones en convertidos 2026", 480), use_container_width=True)

        cols_conv_impacto = [
            "POZO", "BATERIA", "ESTADO_CONVERTIDO_2026", "INTERVENCIONES", "PRCR", "PRAG",
            "PETROLEO_POR_INTERVENCION", "ULTIMO_MES_INTERVENIDO", "PRCR_ULTIMO_MES_INTERVENIDO",
            "POTENCIAL_ULTIMO_MES_INTERVENIDO_BOPD", "POTENCIAL_PROMEDIO_3_MESES_INTERVENCION_BOPD",
            "ULTIMA_FECHA_CON_PRODUCCION_PETROLEO", "PRODUCCION_PETROLEO_ULTIMA_FECHA",
            "FECHA_PRODUCCION_PETROLEO_ANTERIOR", "PRODUCCION_PETROLEO_FECHA_ANTERIOR"
        ]
        cols_conv_impacto = [c for c in cols_conv_impacto if c in conv_det.columns]
        st.dataframe(vista_tabla(conv_det[cols_conv_impacto]), use_container_width=True, hide_index=True)

        st.download_button(
            "Descargar data de gráficos convertidos 2026",
            data=excel_descarga({"Convertidos 2026": conv_det[cols_conv_impacto]}),
            file_name=f"data_convertidos_2026_{anio_obj}_{mes_obj}.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        )

    if not res_clase.empty:
        fig = px.bar(res_clase, x="CLASIFICACION", y=["PRCR", "PRAG"], barmode="group", text_auto=".2f")
        st.plotly_chart(aplicar_layout(fig, "Producción de petróleo y agua por clasificación", 520), use_container_width=True)

    afectados = desp[desp["ESTADO_DESPLAZAMIENTO"].isin(["DEJADO DE HACER", "REDUCIDO"])]

    if not afectados.empty:
        res_afect = afectados.groupby("CLASIFICACION", as_index=False).agg(
            POZOS_AFECTADOS=("POZO_KEY", "nunique"),
            INTERV_BASE_NO_REALIZADAS=("INTERV_BASE_NO_REALIZADAS", "sum"),
            PRCR_BASE_NO_REALIZADO=("PRCR_BASE_NO_REALIZADO", "sum"),
            PRODUCCION_PETROLEO_DEJADA_MES_ANTERIOR=("PRODUCCION_PETROLEO_DEJADA_MES_ANTERIOR", "sum"),
            POTENCIAL_PROMEDIO_3_MESES_INTERVENCION_BOPD=("POTENCIAL_PROMEDIO_3_MESES_INTERVENCION_BOPD", "sum")
        ).sort_values("PRCR_BASE_NO_REALIZADO", ascending=False)

        st.write("Pozos posiblemente desplazados por clasificación")
        st.dataframe(vista_tabla(res_afect), use_container_width=True, hide_index=True)

with tabs[2]:
    st.subheader("Listado de pozos convertidos 2024, 2025 y 2026")

    filtro = st.radio("Año conversión", ["Todos", "2024", "2025", "2026"], horizontal=True)

    tabla = list_conv.copy()
    if filtro != "Todos":
        tabla = tabla[tabla["ANIO_CONVERSION"] == int(filtro)]

    cols = ["ANIO_CONVERSION", "CLASIFICACION", "POZO", "BATERIA", "ESTADO", "INTERVENCIONES", "PRCR", "PRAG", "TIPO_SWAB", "ULTIMA_FECHA_HISTORICA"]
    cols = [c for c in cols if c in tabla.columns]
    st.dataframe(vista_tabla(tabla[cols]), use_container_width=True, hide_index=True)

with tabs[3]:
    st.subheader("Potencial mensual por pozo")

    cols = [
        "POZO", "BATERIA", "CLASIFICACION", "ANIO_CONVERSION",
        "ULTIMO_MES_INTERVENIDO", "ULTIMA_FECHA_INTERVENCION",
        "DIAS_ULTIMO_MES_INTERVENIDO",
        "PRCR_ULTIMO_MES_INTERVENIDO", "PRAG_ULTIMO_MES_INTERVENIDO",
        "INTERV_ULTIMO_MES_INTERVENIDO",
        "POTENCIAL_ULTIMO_MES_INTERVENIDO_BOPD",
        "MES_1_MES_ANTES_ULTIMA_INTERVENCION",
        "PRCR_1_MES_ANTES_ULTIMA_INTERVENCION",
        "POTENCIAL_1_MES_ANTES_ULTIMA_INTERVENCION_BOPD",
        "MES_2_MES_ANTES_ULTIMA_INTERVENCION",
        "PRCR_2_MES_ANTES_ULTIMA_INTERVENCION",
        "POTENCIAL_2_MES_ANTES_ULTIMA_INTERVENCION_BOPD",
        "POTENCIAL_PROMEDIO_3_MESES_INTERVENCION_BOPD",
        "ULTIMA_FECHA_CON_PRODUCCION_PETROLEO",
        "PRODUCCION_PETROLEO_ULTIMA_FECHA",
        "PRODUCCION_AGUA_ULTIMA_FECHA",
        "FECHA_PRODUCCION_PETROLEO_ANTERIOR",
        "PRODUCCION_PETROLEO_FECHA_ANTERIOR",
        "PRODUCCION_AGUA_FECHA_ANTERIOR",
        "ULTIMO_MES_CON_PRODUCCION_PETROLEO",
        "DIAS_ULTIMO_MES_CON_PRODUCCION",
        "PRCR_ULTIMO_MES_CON_PRODUCCION",
        "PRAG_ULTIMO_MES_CON_PRODUCCION",
        "INTERV_ULTIMO_MES_CON_PRODUCCION",
        "POTENCIAL_ULTIMO_MES_CON_PRODUCCION_BOPD"
    ]
    cols = [c for c in cols if c in pot.columns]
    st.dataframe(vista_tabla(pot[cols]), use_container_width=True, hide_index=True)

    top_pot = pot.sort_values("POTENCIAL_PROMEDIO_3_MESES_INTERVENCION_BOPD", ascending=False).head(top_n)
    if not top_pot.empty:
        fig = px.bar(
            top_pot,
            x="POZO",
            y=["POTENCIAL_ULTIMO_MES_INTERVENIDO_BOPD", "POTENCIAL_1_MES_ANTES_ULTIMA_INTERVENCION_BOPD", "POTENCIAL_2_MES_ANTES_ULTIMA_INTERVENCION_BOPD"],
            barmode="group",
            hover_data=["BATERIA", "CLASIFICACION", "ULTIMO_MES_INTERVENIDO", "MES_1_MES_ANTES_ULTIMA_INTERVENCION", "MES_2_MES_ANTES_ULTIMA_INTERVENCION", "ULTIMA_FECHA_CON_PRODUCCION_PETROLEO", "ULTIMO_MES_CON_PRODUCCION_PETROLEO"]
        )
        st.plotly_chart(aplicar_layout(fig, f"Top {top_n} pozos por potencial", 560), use_container_width=True)

with tabs[4]:
    st.subheader("Pozos y baterías")

    st.dataframe(vista_tabla(res_pozos), use_container_width=True, hide_index=True)

    if not res_bat.empty:
        fig = px.bar(
            res_bat.head(top_n),
            x="BATERIA",
            y="PRCR",
            text="PRCR",
            hover_data=["POZOS_INTERVENIDOS", "POZOS_NO_INTERVENIDOS", "INTERVENCIONES", "PRAG"]
        )
        st.plotly_chart(aplicar_layout(fig, f"Top {top_n} baterías por producción de petróleo", 520), use_container_width=True)

    if not res_tend.empty:
        orden = [MESES[m] for m in sorted(res_tend["MES"].unique())]
        fig2 = px.line(
            res_tend,
            x="MES_NOMBRE",
            y=["PRCR", "PRAG"],
            markers=True,
            category_orders={"MES_NOMBRE": orden}
        )
        st.plotly_chart(aplicar_layout(fig2, f"Tendencia mensual {anio_obj}", 520), use_container_width=True)

with tabs[5]:
    st.subheader("TS y CS")

    if res_tipo.empty:
        st.info("Sin información por tipo de swab.")
    else:
        col1, col2 = st.columns(2)
        with col1:
            fig = px.bar(res_tipo, x="TIPO_SWAB", y=["PRCR", "PRAG"], barmode="group", text_auto=".2f")
            st.plotly_chart(aplicar_layout(fig, "Producción de petróleo y agua por tipo de swab", 520), use_container_width=True)
        with col2:
            figp = px.pie(res_tipo, names="TIPO_SWAB", values="INTERVENCIONES", hole=0.45)
            figp.update_traces(texttemplate="%{percent:.2%}")
            st.plotly_chart(aplicar_layout(figp, "Distribución de intervenciones por tipo de swab", 520), use_container_width=True)

        st.dataframe(vista_tabla(res_tipo), use_container_width=True, hide_index=True)

with tabs[6]:
    st.subheader("Descargas")

    tablas = {
        "Resumen pozos": res_pozos,
        "Analisis pozos dejados": desp,
        "Listado convertidos": list_conv,
        "Potencial pozos": pot,
        "Resumen clasificacion": res_clase,
        "Resumen baterias": res_bat,
        "Resumen TS CS": res_tipo,
        "Tendencia mensual": res_tend,
        "Pozos ATA excluidos": df_ata,
        "Data graf dejados": (
            desp[desp["ESTADO_DESPLAZAMIENTO"].isin(["DEJADO DE HACER", "REDUCIDO"])]
            .sort_values("PRCR_BASE_NO_REALIZADO", ascending=False)
            .head(top_n)
        ),
        "Data graf potencial": pot.sort_values("POTENCIAL_PROMEDIO_3_MESES_INTERVENCION_BOPD", ascending=False).head(top_n)
    }

    st.download_button(
        "Descargar toda la información en Excel",
        data=excel_descarga(tablas),
        file_name=f"swab_lote_x_{anio_obj}_{mes_obj}.xlsx",
        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    )

    kpis_ppt = {
        "pozos_universo": pozos_universo,
        "pozos_interv": pozos_interv,
        "pozos_no": pozos_no,
        "interv": interv,
        "petroleo": petroleo,
        "agua": agua
    }

    ppt_bytes = crear_ppt_dashboard(
        periodo=periodo,
        kpis=kpis_ppt,
        res_clase=res_clase,
        res_bat=res_bat,
        res_tipo=res_tipo,
        res_tend=res_tend,
        desp=desp,
        pot=pot,
        list_conv=list_conv,
        actual=actual,
        top_n=top_n
    )

    st.download_button(
        "Descargar PPT editable",
        data=ppt_bytes,
        file_name=f"dashboard_swab_lote_x_{anio_obj}_{mes_obj}.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )

    validacion = pd.DataFrame({
        "Concepto": [
            "Hoja usada", "Fecha mínima", "Fecha máxima", "Registros analizados",
            "Pozos SWAB analizados", "Pozos ATA excluidos", "Registros ATA excluidos",
            "Convertidos 2026 fijos", "PRCR", "PRAG"
        ],
        "Valor": [
            hoja_usada, str(fecha_min), str(fecha_max), f"{len(df):,}",
            f"{df['POZO_KEY'].nunique():,}", f"{df_ata['POZO_KEY'].nunique():,} de 88",
            f"{len(df_ata):,}", "20", "Producción de petróleo", "Producción de agua"
        ]
    })

    st.write("Validación")
    st.dataframe(vista_tabla(validacion), use_container_width=True, hide_index=True)
